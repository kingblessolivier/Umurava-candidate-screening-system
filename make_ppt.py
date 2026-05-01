from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color Palette (minimal: 3 colors only) ──────────────────────────────────
DARK   = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
BLUE   = RGBColor(0x38, 0x8B, 0xF5)   # single accent blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)   # slide background
MGRAY  = RGBColor(0x64, 0x74, 0x8B)   # secondary text
DGRAY  = RGBColor(0x1E, 0x29, 0x3B)   # body text

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape

def add_text(slide, text, x, y, w, h, size, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_para(tf, text, size, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def slide_bg(slide, color=LGRAY):
    add_rect(slide, 0, 0, W, H, color)

def header_bar(slide, title, subtitle=None):
    """Dark top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.35), DARK)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.82), Inches(11),
                 Inches(0.42), 14, color=BLUE)

def accent_line(slide, y=Inches(1.35)):
    add_rect(slide, 0, y, W, Pt(3), BLUE)

def bullet_card(slide, x, y, w, h, title, bullets, icon=""):
    """White card with title and bullet list."""
    card = add_rect(slide, x, y, w, h, WHITE)
    card.shadow.inherit = False
    add_rect(slide, x, y, w, Pt(4), BLUE)
    tx = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12),
                                  w - Inches(0.25), h - Inches(0.18))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (icon + "  " if icon else "") + title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = "Calibri"
    for b in bullets:
        add_para(tf, "  • " + b, 11, color=DGRAY, space_before=3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, Inches(3.1), W, Pt(3), BLUE)

add_text(s, "TalentAI", Inches(0.8), Inches(1.5), Inches(11), Inches(1.6),
         72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "AI-Powered Talent Screening System",
         Inches(0.8), Inches(3.25), Inches(11.7), Inches(0.7),
         22, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Umurava AI Hackathon  ·  2025",
         Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.5),
         14, color=MGRAY, align=PP_ALIGN.CENTER)
add_text(s, "Next.js  ·  Node.js  ·  MongoDB  ·  Google Gemini 2.0 Flash",
         Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
         11, color=MGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "The Problem", "Why talent screening is broken today")
accent_line(s)

problems = [
    ("📋  Manual & Slow",
     ["Recruiters spend 23 hrs/week just on screening",
      "Most resumes receive < 10 seconds of attention"]),
    ("⚖️  Biased by Default",
     ["Unconscious bias filters out qualified candidates",
      "Gender language & prestige bias are invisible to humans"]),
    ("🔍  No Explainability",
     ["Candidates never learn *why* they were rejected",
      "Hiring decisions can't be audited or improved"]),
    ("📊  One-Size-Fits-All",
     ["Every job requires different skill weightings",
      "Generic ATS tools ignore role-specific priorities"]),
]

cols = [(Inches(0.4),  Inches(1.55)),
        (Inches(6.85), Inches(1.55)),
        (Inches(0.4),  Inches(4.2)),
        (Inches(6.85), Inches(4.2))]

for (px, py), (title, bullets) in zip(cols, problems):
    bullet_card(s, px, py, Inches(6.2), Inches(2.35), title, bullets)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Our Solution — TalentAI", "One platform. End-to-end AI screening.")
accent_line(s)

add_text(s,
    "TalentAI automates the entire candidate screening pipeline — from raw resume PDF to "
    "a ranked shortlist with confidence scores, bias alerts, and personalised rejection "
    "feedback — powered by Google Gemini 2.0 Flash.",
    Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.0), 14, color=DGRAY)

stats = [
    ("3-Stage", "Gemini AI Pipeline"),
    ("0–100",   "Confidence Score per candidate"),
    ("5",       "Live Analytics Charts"),
    ("25×",     "Batch processing per API call"),
    ("100%",    "Explainable rejections"),
]

bx = Inches(0.4)
for i, (num, label) in enumerate(stats):
    cx = bx + i * Inches(2.5)
    add_rect(s, cx, Inches(2.75), Inches(2.28), Inches(1.65), WHITE)
    add_rect(s, cx, Inches(2.75), Inches(2.28), Pt(4), BLUE)
    add_text(s, num, cx, Inches(2.88), Inches(2.28), Inches(0.75),
             30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, label, cx, Inches(3.6), Inches(2.28), Inches(0.7),
             11, color=DGRAY, align=PP_ALIGN.CENTER)

features = [
    "Resume PDF parsing → structured TalentProfile stored in MongoDB",
    "AI Enhance Job: rough description → structured requirements + scoring weights",
    "Custom weight sliders per job (Skills / Experience / Education / Projects / Availability)",
    "Risk detection: unverified skills, short tenures, employment gaps",
    "Bias detection: gender language, age indicators, institution prestige bias",
]
tx = slide.shapes.add_textbox if False else s.shapes.add_textbox(
    Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "What's included:"
run.font.size = Pt(12)
run.font.bold = True
run.font.color.rgb = DARK
run.font.name = "Calibri"
for f in features:
    add_para(tf, "  ✓  " + f, 11, color=DGRAY, space_before=2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Tech Stack", "Modern, production-grade architecture")
accent_line(s)

stack = [
    ("Frontend", [
        "Next.js 15  (App Router + TypeScript)",
        "Tailwind CSS  •  Redux Toolkit",
        "Recharts  (5 analytics charts)",
        "Dark mode support  •  PDF preview",
    ]),
    ("Backend", [
        "Node.js + Express + TypeScript",
        "MongoDB + Mongoose",
        "Google Gemini 2.0 Flash (AI engine)",
        "Multer  (PDF upload)  •  JWT Auth",
    ]),
    ("AI Layer", [
        "Preprocessing service → numeric signals",
        "Prompt 1: Job Understanding",
        "Prompt 2: Multi-Candidate Evaluation",
        "Prompt 3: Final Ranking + Why Not Selected",
    ]),
    ("Infrastructure", [
        "RESTful API  •  Batch processing (25/call)",
        "Multi-batch re-ranking for 50+ candidates",
        "Resume upsert  (update existing profiles)",
        "Confidence scoring  (0–100 per candidate)",
    ]),
]

cols2 = [(Inches(0.4),  Inches(1.6)),
         (Inches(3.65), Inches(1.6)),
         (Inches(6.9),  Inches(1.6)),
         (Inches(10.15),Inches(1.6))]

for (px, py), (title, bullets) in zip(cols2, stack):
    bullet_card(s, px, py, Inches(3.0), Inches(5.5), title, bullets)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI Pipeline
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "3-Stage Gemini AI Pipeline", "How TalentAI thinks about your candidates")
accent_line(s)

stages = [
    ("Stage 1", "Job Understanding",
     "Parse raw job description → extract must-have skills, nice-to-have "
     "skills, experience range, education requirements, and apply custom "
     "scoring weights.  AI Enhance auto-generates requirements from a rough brief."),
    ("Stage 2", "Multi-Candidate Evaluation",
     "Batch-evaluate 25 candidates per API call against the parsed job spec. "
     "Preprocessing service extracts numeric signals (GPA, years of experience, "
     "tenure lengths) first to reduce Gemini hallucination."),
    ("Stage 3", "Final Ranking + Explanations",
     "Re-rank across batches, assign 0–100 Confidence Score, flag bias signals, "
     "detect risks, and generate personalised 'Why Not Selected' feedback with "
     "actionable improvement suggestions for every rejected candidate."),
]

for i, (stage, title, desc) in enumerate(stages):
    y = Inches(1.65) + i * Inches(1.82)
    add_rect(s, Inches(0.4), y, Inches(1.5), Inches(1.55), DARK)
    add_text(s, stage, Inches(0.4), y + Inches(0.15), Inches(1.5), Inches(0.45),
             11, color=MGRAY, align=PP_ALIGN.CENTER)
    add_text(s, str(i + 1), Inches(0.4), y + Inches(0.45), Inches(1.5), Inches(0.75),
             42, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(2.1), y, Inches(10.8), Inches(1.55), WHITE)
    add_rect(s, Inches(2.1), y, Inches(10.8), Pt(4), BLUE)
    tx2 = s.shapes.add_textbox(Inches(2.28), y + Inches(0.1),
                                Inches(10.5), Inches(1.35))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = title
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = DARK
    run2.font.name = "Calibri"
    add_para(tf2, desc, 11, color=DGRAY, space_before=4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Differentiators
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Key Differentiators", "What makes TalentAI stand apart")
accent_line(s)

diffs = [
    ("🎯  Confidence Score",
     ["0–100 numeric score per candidate",
      "Based on verified signals, not vibes",
      "Fully explainable to hiring managers"]),
    ("⚖️  Bias Detection",
     ["Flags gender-coded language",
      "Detects age indicators",
      "Warns on institution prestige bias"]),
    ("💬  Why Not Selected?",
     ["Every rejected candidate gets feedback",
      "Constructive improvement suggestions",
      "Auditable, fair, legally defensible"]),
    ("⚙️  Custom Weights",
     ["Per-job sliders: Skills / Exp / Edu",
      "Projects & Availability configurable",
      "Weights must sum to 100"]),
    ("🔍  Risk Detection",
     ["Unverified skill claims flagged",
      "Short tenures & gaps highlighted",
      "Recruiter attention drawn automatically"]),
    ("📊  Analytics Dashboard",
     ["Top skills bar chart",
      "Skill gap analysis",
      "Candidate source & recommendation pies",
      "Top vs Avg radar chart"]),
]

positions = [
    (Inches(0.4),  Inches(1.6)),
    (Inches(4.55), Inches(1.6)),
    (Inches(8.7),  Inches(1.6)),
    (Inches(0.4),  Inches(4.15)),
    (Inches(4.55), Inches(4.15)),
    (Inches(8.7),  Inches(4.15)),
]
for (px, py), (title, bullets) in zip(positions, diffs):
    bullet_card(s, px, py, Inches(3.9), Inches(2.2), title, bullets)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Analytics Dashboard
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Analytics Dashboard", "Data-driven insights at a glance")
accent_line(s)

charts = [
    ("Bar Chart", "Top Skills in Pool",
     "See which skills dominate your applicant pool vs job requirements"),
    ("Bar Chart", "Skill Gap Analysis",
     "Identify where candidates fall short of the job spec"),
    ("Pie Chart", "Recommendation Breakdown",
     "Strong Yes / Yes / Maybe / No distribution at a glance"),
    ("Pie Chart", "Candidate Sources",
     "Track where your best candidates are coming from"),
    ("Radar Chart", "Top Candidate vs Average",
     "Visual comparison of your best pick against the applicant average"),
]

for i, (ctype, ctitle, cdesc) in enumerate(charts):
    y = Inches(1.6) + i * Inches(1.06)
    add_rect(s, Inches(0.4), y, Inches(1.6), Inches(0.9), DARK)
    add_text(s, ctype, Inches(0.4), y + Inches(0.05), Inches(1.6), Inches(0.8),
             10, color=BLUE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(2.15), y, Inches(10.75), Inches(0.9), WHITE)
    add_rect(s, Inches(2.15), y, Inches(10.75), Pt(3), BLUE)
    tx3 = s.shapes.add_textbox(Inches(2.35), y + Inches(0.05),
                                Inches(10.4), Inches(0.8))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = ctitle
    r3.font.size = Pt(13)
    r3.font.bold = True
    r3.font.color.rgb = DARK
    r3.font.name = "Calibri"
    add_para(tf3, cdesc, 11, color=MGRAY, space_before=2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Demo Flow
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Live Demo Flow", "End-to-end walkthrough")
accent_line(s)

steps = [
    ("1", "Create Job",       "Post a new role + rough description  →  AI Enhance fills requirements & sets weights"),
    ("2", "Upload Resumes",   "Drag-and-drop PDF batch  →  Gemini parses each into a structured TalentProfile"),
    ("3", "Run Screening",    "Hit 'Screen Candidates'  →  3-stage AI pipeline evaluates, ranks, explains"),
    ("4", "Review Results",   "Sorted shortlist with scores, risk flags, bias alerts, and 'Why Not Selected'"),
    ("5", "Analytics",        "Dashboard auto-populates with 5 live charts showing pool-level insights"),
]

for i, (num, stitle, sdesc) in enumerate(steps):
    y = Inches(1.6) + i * Inches(1.06)
    add_rect(s, Inches(0.4), y, Inches(0.65), Inches(0.9), BLUE)
    add_text(s, num, Inches(0.4), y + Inches(0.06), Inches(0.65), Inches(0.75),
             22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(1.2), y, Inches(11.7), Inches(0.9), WHITE)
    tx4 = s.shapes.add_textbox(Inches(1.4), y + Inches(0.06),
                                Inches(11.4), Inches(0.8))
    tf4 = tx4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.LEFT
    r4a = p4.add_run()
    r4a.text = stitle + "   "
    r4a.font.size = Pt(13)
    r4a.font.bold = True
    r4a.font.color.rgb = DARK
    r4a.font.name = "Calibri"
    r4b = p4.add_run()
    r4b.text = "— " + sdesc
    r4b.font.size = Pt(11)
    r4b.font.bold = False
    r4b.font.color.rgb = MGRAY
    r4b.font.name = "Calibri"

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Thank You
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, Inches(3.55), W, Pt(3), BLUE)

add_text(s, "Thank You", Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.4),
         62, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Questions?  Let's talk.",
         Inches(0.8), Inches(2.75), Inches(11.7), Inches(0.65),
         20, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "nsengimana_222023259@stud.ur.ac.rw",
         Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.5),
         13, color=MGRAY, align=PP_ALIGN.CENTER)
add_text(s, "TalentAI  ·  Umurava AI Hackathon 2025",
         Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.45),
         12, color=MGRAY, align=PP_ALIGN.CENTER)

add_text(s,
    "AI-Powered  ·  Explainable  ·  Bias-Aware  ·  Production-Ready",
    Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.45),
    11, color=MGRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\USER\Documents\projects\TalentAi\TalentAI_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
